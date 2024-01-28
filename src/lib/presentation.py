import os
import re
from pptx import Presentation
from pptx import Presentation
from pptx.util import Pt, Mm, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, PP_PLACEHOLDER
from PIL import Image

from lib.config_manager import get_processed_slides_from_json_file, get_slide_id, image_basename, load_config, read_config_from_text_box, write_config_to_text_box, write_default_json_config_to_text_frame, write_processed_slide_to_json_file
from pptx.enum.shapes import MSO_SHAPE
import shutil
from datetime import datetime


config = load_config()
presentation_filename = config['presentation_filename']
if not presentation_filename.endswith('.pptx'):
    presentation_filename += '.pptx'
header_title = config['header_title']
default_comment = config['default_comment']
page_string = config['page_string']
header_title = config['header_title']
color_code_sender = config['color_code_sender']

page_string_regex_pattern = page_string.replace("%d", r"\d+")
page_string_regex_pattern = "^" + page_string_regex_pattern
regex_right_header = re.compile(page_string_regex_pattern)

header_title_regex_pattern = "^" + header_title
regex_left_header = re.compile(header_title_regex_pattern)

slides_dict = {}
# regex_right_header = re.compile(r"^Seite \d+ von \d+")  # r"Seite \d+/\d+"
# regex_right_header = re.compile(
#     r"^Stellungnahme Nr. \d+ - Seite \d+ von \d+")  # r"Seite \d+/\d+"
# regex_left_header = re.compile(
#     r"^Abwägungsvorschlag Träger öffentlicher Belange")


# Create a backup directory if it doesn't exist
backup_dir = "backup"
os.makedirs(backup_dir, exist_ok=True)

# Check if the PowerPoint file already exists
if os.path.exists(presentation_filename):
    # If it does, create a backup
    base_name, ext = os.path.splitext(os.path.basename(presentation_filename))
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")  # Current timestamp
    backup_filename = os.path.join(
        backup_dir, f"{base_name}.backup.{timestamp}{ext}")
    shutil.copy(presentation_filename, backup_filename)

    # Then open the original file
    prs = Presentation(presentation_filename)
    starting_page_count = len(prs.slides)

else:
    prs = Presentation()

prs.slide_width = Mm(297)
prs.slide_height = Mm(210)


def is_duplicate(image, sender):
    # Check if the image is already in the presentation file or in the slides_dict. slides_dict is used to keep track of images that are not in the presentation file yet.
    return any((read_config_from_text_box(slide, "image") == image
                and read_config_from_text_box(slide, "sender") == sender) for slide in prs.slides) or any(image in images for images in slides_dict.values())


def add_header_left(slide):
    bold = False
    size = 11
    # Add a header to the slide
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = header_title
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_header_right(slide, page_number, total_pages, sender_nr=0):
    bold = False
    size = 10
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = page_string % (sender_nr, page_number, total_pages)
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_image(slide, image_filepath):
    # Define the maximum height and width
    max_height = Cm(18)
    max_width = Cm(14)

    # Open the image and get its size
    image = Image.open(image_filepath)
    image_width, image_height = image.size
    aspect_ratio = image_width / image_height

    # Calculate the height and width, maintaining the aspect ratio
    width_by_height = max_height * aspect_ratio
    height_by_width = max_width / aspect_ratio

    if width_by_height <= max_width:
        # Image fits within max dimensions by adjusting height
        width = width_by_height
        height = max_height
    else:
        # Image fits within max dimensions by adjusting width
        width = max_width
        height = height_by_width

    # Add the image to the slide
    top = Cm(1.6)
    left = Cm(0.5)
    pic = slide.shapes.add_picture(
        image_filepath, left, top, height=height, width=width)


def add_source_text(slide, text):
    # Function that adds a string to the top left of the slide
    left = Cm(1.6)
    top = Cm(20)
    width = Cm(13.5)
    height = Cm(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(6)


def add_divider_line(slide, prs):
    left = prs.slide_width // 2
    top = Cm(0.7)
    end_top = prs.slide_height
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, left, top, left, end_top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False

    # Add horizontal divider line
    left = Cm(0.5)
    top = Cm(0.7)
    end_left = prs.slide_width - Cm(0.5)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, left, top, end_left, top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False


def add_text_box(slide):
    left = Cm(16)
    top = Cm(1)
    width = Cm(11)
    height = Cm(20)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = default_comment
    p.font.bold = False
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.LEFT


def remove_title_placeholder(slide):
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == 0 and phf.type == PP_PLACEHOLDER.TITLE:
                sp = shape._element
                sp.getparent().remove(sp)


def add_image_to_presentation(image, sender):
    slide_id = get_slide_id(image, sender)
    if slide_id not in get_processed_slides_from_json_file():
        if sender in slides_dict:
            # Add the image to the sender's list
            slides_dict[sender].append(image)
        else:
            # If the sender does not exist, create a new list for them
            slides_dict[sender] = [image]
        write_processed_slide_to_json_file(slide_id)
    else:
        print(f'Slide was removed manually. Not re-adding. {slide_id}')


def create_presentation_from_dict():
    i = 0
    for sender, images in slides_dict.items():
        for image in images:
            i += 1
            # Add a new slide at the end
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            write_default_json_config_to_text_frame(slide)
            write_config_to_text_box(
                slide, "image", image)
            write_config_to_text_box(slide, "sender", sender)
            remove_title_placeholder(slide)
            add_image(slide, image)
            text = f'Quelle: {sender} - {image_basename(image)}'
            if len(text) > 120:
                space_index = text.rfind(' ', 0, 120)
                if space_index != -1:
                    text = text[:space_index] + "\n" + text[space_index+1:]
            add_source_text(slide, text)
            add_text_box(slide)
            add_divider_line(slide, prs)
    return prs


def add_left_border(slide, color=RGBColor(255, 255, 0)):
    # Create a new shape (line) on the slide
    slide_height = Mm(210)
    left_border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Pt(5), slide_height
    )

    # Set the fill of the shape to yellow
    fill = left_border.fill
    fill.solid()
    fill.fore_color.rgb = color

    # Remove the outline of the shape
    left_border.line.fill.background()


def remove_left_border(slide):
    for shape in slide.shapes:
        # Check if the shape is a rectangle
        if shape.shape_type == MSO_SHAPE.RECTANGLE:
            # Check if the shape is on the left side of the slide
            if shape.left == 0:
                # Remove the shape
                slide.shapes._spTree.remove(shape._element)


def add_headers(prs):
    page_count = len(prs.slides)
    sender_count = 0
    color_switch = True  # Variable to switch colors
    current_sender = None  # Variable to keep track of the current sender

    for i, slide in enumerate(prs.slides):
        remove_left_border(slide)
        # if right header exists, remove it
        sender = read_config_from_text_box(slide, "sender")

        # Check if the sender has changed
        if sender != current_sender:
            color_switch = not color_switch  # Switch color
            current_sender = sender  # Update the current sender
            sender_count += 1  # Increment the sender count

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if regex_right_header.match(text_frame.text):
                    # print("removing right header")
                    slide.shapes._spTree.remove(shape._element)
                if regex_left_header.match(text_frame.text):
                    # print("removing left header")
                    slide.shapes._spTree.remove(shape._element)

        if color_code_sender:
            # Add header with alternating colors
            color = RGBColor(0, 0, 255) if color_switch else RGBColor(
                255, 255, 0)  # Blue if color_switch is True, otherwise Yellow
            add_left_border(slide, color)

            add_header_right(slide, i+1, page_count, sender_count)
            add_header_left(slide)
    return page_count


def save_presentation(prs):
    prs.save(presentation_filename)
