import os
import re
import json
import yaml
import shutil
from html import escape
from email import policy
from email.parser import BytesParser
from email.utils import parsedate_to_datetime, parseaddr
from weasyprint import HTML
from pdf2image import convert_from_path
from PIL import Image, ImageChops
from pptx import Presentation
from pptx.util import Pt, Mm, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, PP_PLACEHOLDER
from mailparser_reply import EmailReplyParser

# conda install conda-forge::weasyprint

# if config.yaml does not exist, create it
if not os.path.exists('config.yaml'):
    print("config.yaml does not exist. Creating it with default values.")
    with open('config.yaml', 'w') as f:
        f.write('presentation_filename: presentation.pptx\n')
        f.write('header_title: "Eingegangene Emails"\n')
        f.write('default_comment: "Wird zur Kenntnis genommen."\n')
        f.write('pdf_blacklist:\n')
        f.write('  - ".*DUMMY.*"\n')
        f.write('eml_input_dir: eml\n')
        f.write('scanned_input_dir: scanned\n')
        f.write('output_dir: output\n')
        # f.write('font_path: "C:\\Windows\\Fonts\\Arial.ttf"\n')
        f.write('font_path: "/Library/Fonts/Arial.ttf"\n')
        f.close()
    print("Please configure config.yaml and run the script again.")
    exit()

with open('config.yaml', 'r') as f:
    config = yaml.safe_load(f)

presentation_filename = config['presentation_filename']
header_title = config['header_title']
default_comment: config['default_comment']
pdf_blacklist = [re.compile(pattern) for pattern in config['pdf_blacklist']]
eml_input_dir = config['eml_input_dir']
scanned_input_dir = config['scanned_input_dir']
output_dir = config['output_dir']
font_path = config['font_path']

# the order of these types is important. The first one will be used if both are present
# when preferring html, the issue is that reply/response deztect is bad (nil even?)
text_content_types = ['text/plain', 'text/html']
# text_content_types = ['text/html', 'text/plain']
starting_page_count = 0
slides_dict = {}
# regex to match the right header
regex_right_header = re.compile(r"Seite \d+/\d+")

# Create output directory if it doesn't exist
if not os.path.exists(output_dir):
    os.makedirs(output_dir)


# Check if the PowerPoint file already exists
if os.path.exists(presentation_filename):
    # If it does, open it
    prs = Presentation(presentation_filename)
    starting_page_count = len(prs.slides)

else:
    prs = Presentation()

# Set slide width and height to A4 portrait dimensions (210mm x 297mm)
prs.slide_width = Mm(297)
prs.slide_height = Mm(210)


def image_basename(image):
    return os.path.splitext(os.path.basename(image))[0]


def get_image_name(image):
    return image


def get_sender_name(sender):
    return sender


def write_default_json_config_to_text_frame(slide):
    # create a new text box, not a note
    text_frame = slide.shapes.add_textbox(Cm(0.1), Cm(0), Cm(0.1), Cm(
        0.1)).text_frame  # Parameters are left, top, width, height
    # write default json to it
    default_for_json = {"image_name": "default",
                        "hashed_sender_name": "default"}
    json_hashed_image_name = json.dumps(default_for_json)
    text_frame.text = json_hashed_image_name
    text_frame.paragraphs[0].font.size = Pt(1)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def get_json_text_box(slide):
    # search all text boxes and find the first one that contains json
    regex_json = re.compile(r"^{.*}$")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if regex_json.match(text_frame.text):
                return text_frame
    return None


def write_config_to_text_frame(slide, key, value):
    json_text_box = get_json_text_box(slide)
    if json_text_box is None:
        return
    # get the json from the text box and update it
    existing_json = json.loads(json_text_box.text)
    existing_json[key] = value
    # write the updated json back to the text box
    json_text_box.text = json.dumps(existing_json)
    json_text_box.paragraphs[0].font.size = Pt(1)
    json_text_box.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def read_config_from_text_frame(slide, key):
    json_text_box = get_json_text_box(slide)
    if json_text_box is None:
        return None
    # get the json from the text box
    existing_json = json.loads(json_text_box.text)
    return existing_json.get(key, None)


def sanitize_filename(filename):
    return re.sub(r'[\\/*?:"<>|\[\]]', '_', filename)


def remove_quoted(email_message):
    email_message = EmailReplyParser(
        languages=['de', 'en']).parse_reply(text=email_message)

    # lines = email_message.split('\n')
    # non_quoted_lines = []
    # for line in lines:
    #     if line.startswith(('Von:', 'Gesendet:', 'An:', 'Betreff:')):
    #         break
    #     non_quoted_lines.append(line)
    # email_message = '\n'.join(non_quoted_lines)

    return email_message


def get_html_content(msg):
    for type in text_content_types:
        if msg.is_multipart():
            for part in msg.iter_parts():
                if part.get_content_type() == type:
                    # if Content-Disposition is attachment, skip
                    if part.get('Content-Disposition') and part.get('Content-Disposition') != "inline":
                        continue
                    reply = remove_quoted(part.get_content())
                    if reply:
                        return {"type": type, "content": reply}
                    return {"type": type, "content": part.get_content()}
                elif part.is_multipart():
                    return get_html_content(part)
        else:
            if msg.get_content_type() == type:
                if msg.get('Content-Disposition') and msg.get('Content-Disposition') != "inline":
                    continue
                reply = remove_quoted(msg.get_content())
                if reply:
                    return {"type": type, "content": reply}
                return {"type": type, "content": msg.get_content()}

    return None


def add_header_left(slide):
    bold = False
    size = 11
    # Add a header to the slide
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = header_title
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_header_right(slide, page_number, total_pages):
    bold = False
    size = 10
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = f'Seite {page_number}/{total_pages}'
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def crop_whitespace(image):
    # Crop the image to remove whitespace
    bg = Image.new(image.mode, image.size, image.getpixel((0, 0)))
    diff = ImageChops.difference(image, bg)
    diff = ImageChops.add(diff, diff, 2.0, -100)
    bbox = diff.getbbox()
    if bbox:
        # Modify the bounding box to keep the full width of the image
        bbox = (0, bbox[1], image.size[0], bbox[3])
        return image.crop(bbox)


def add_image(slide, image_filepath):
    # Define the maximum height and width
    max_height = Cm(18)
    max_width = Cm(14)

    # Open the image and get its size
    image = Image.open(image_filepath)
    image_width, image_height = image.size
    aspect_ratio = image_width / image_height

    # Calculate the height and width, maintaining the aspect ratio
    width_by_height = max_height * aspect_ratio
    height_by_width = max_width / aspect_ratio

    if width_by_height <= max_width:
        # Image fits within max dimensions by adjusting height
        width = width_by_height
        height = max_height
    else:
        # Image fits within max dimensions by adjusting width
        width = max_width
        height = height_by_width

    # Add the image to the slide
    top = Cm(1.6)
    left = Cm(0.5)
    pic = slide.shapes.add_picture(
        image_filepath, left, top, height=height, width=width)


def add_source_text(slide, text):
    # Function that adds a string to the top left of the slide
    left = Cm(1.6)
    top = Cm(20)
    width = Cm(13.5)
    height = Cm(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(6)


def add_divider_line(slide, prs):
    left = prs.slide_width // 2
    top = Cm(0.7)
    end_top = prs.slide_height
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, left, top, left, end_top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False

    # Add horizontal divider line
    left = Cm(0.5)
    top = Cm(0.7)
    end_left = prs.slide_width - Cm(0.5)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, left, top, end_left, top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False


def add_text_box(slide):
    left = Cm(16)
    top = Cm(1)
    width = Cm(11)
    height = Cm(20)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Zur Kenntnis genommen."
    p.font.bold = False
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.LEFT


def remove_title_placeholder(slide):
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == 0 and phf.type == PP_PLACEHOLDER.TITLE:
                sp = shape._element
                sp.getparent().remove(sp)


def create_directories(output_dir, sender):
    # Create a new directory with the same name as the sender's email address
    sender_dir = f'{output_dir}/{sender}'
    os.makedirs(sender_dir, exist_ok=True)

    # Create a subdirectory called "attachments" to store attachments
    attachments_dir = f'{sender_dir}/attachments'
    os.makedirs(attachments_dir, exist_ok=True)

    # Create a subdirectory called "text" to store HTML files
    text_dir = f'{sender_dir}/text'
    os.makedirs(text_dir, exist_ok=True)

    # Create a subdirectory called "img" in subdirectory "attachments" to store images
    attachments_img_dir = f'{attachments_dir}/img'
    os.makedirs(attachments_img_dir, exist_ok=True)

    # Create a subdirectory called "img" in subdirectory "text" to store images
    text_img_dir = f'{text_dir}/img'
    os.makedirs(text_img_dir, exist_ok=True)

    return sender_dir, attachments_dir, text_dir, attachments_img_dir, text_img_dir


def get_email_date(msg):
    email_date = parsedate_to_datetime(msg['Date'])
    # Format the date and time as YYYY-MM-DD_HH_MM_SS
    formatted_date_time = email_date.strftime('%Y-%m-%d_%H_%M_%S')
    return formatted_date_time


def handle_pdf(part, msg, output_dir):
    payload = part.get_payload(decode=True)
    if payload is not None:
        attachment_filename = part.get_filename()
        if is_in_blacklist(attachment_filename):
            return None
        formatted_date = get_email_date(msg)
        filename, file_extension = os.path.splitext(attachment_filename)
        filename_with_date = f"{filename}_{formatted_date}{file_extension}"
        filepath = os.path.join(output_dir, filename_with_date)
        with open(filepath, 'wb') as f:
            f.write(payload)
        return filepath


def handle_zip(part, msg, output_dir):
    # TODO: handle zip files. unpack them, save the files and add them to the presentation
    payload = part.get_payload(decode=True)
    if payload is not None:
        attachment_filename = part.get_filename()
        if is_in_blacklist(attachment_filename):
            return None
        filepath = os.path.join(output_dir, attachment_filename)
        print(
            f'Found zip file. Saving to: {attachment_filename} in path {filepath}')
        with open(filepath, 'wb') as f:
            f.write(payload)


def handle_other_attachments(part, msg, output_dir):
    payload = part.get_payload(decode=True)
    if payload is not None:
        attachment_filename = part.get_filename()
        if is_in_blacklist(attachment_filename):
            return None
        filepath = os.path.join(output_dir, attachment_filename)
        if attachment_filename.endswith('.pdf'):
            print(
                f'Found PDF attachment without application/pdf MIME type. Saving to: {attachment_filename} in path {filepath}')
            return handle_pdf(part, msg, output_dir)
        print(
            f'Found non-PDF attachment. Saving to: {attachment_filename} in path {filepath}')
        with open(filepath, 'wb') as f:
            f.write(payload)


def is_in_blacklist(filename):
    for regex in pdf_blacklist:
        if regex.match(filename):
            print(f'Found blacklisted file: {filename}')
            return True
    return False


def save_attachments(msg, output_dir):
    filepaths = []
    if msg.is_multipart():
        for part in msg.iter_parts():
            if part.get_content_type() == 'application/pdf':
                fp = handle_pdf(part, msg, output_dir)
                if fp is not None:
                    filepaths.append(fp)
            elif part.get_content_type() == 'application/zip':
                handle_zip(part, msg, output_dir)
            elif part.get('Content-Disposition') and part.get_filename() is not None:
                fp = handle_other_attachments(part, msg, output_dir)
                if fp is not None:
                    filepaths.append(fp)
    return filepaths


def convert_pdfs_to_images(pdf_filepaths, output_dir):
    images = []
    for pdf_filepath in pdf_filepaths:
        # Convert the PDF to images
        pdf_images = convert_from_path(pdf_filepath)
        # Save each image under the output directory
        for i, pdf_image in enumerate(pdf_images):
            pdf_image = crop_whitespace(pdf_image)
            image_filepath = f'{output_dir}/{os.path.basename(pdf_filepath)}_{i}.png'
            if pdf_image is not None:
                pdf_image.save(image_filepath, 'PNG')
                images.append(image_filepath)
    return images


def create_pdf(message_content, output_dir, filename):
    content = message_content["content"]
    type = message_content["type"]
    # Check if the content is HTML or plain text
    if type == "text/plain":
        # The content is plain text
        # Set a maximum width to fit the content within an A4 page
        content = escape(content).replace('\n', '<br>')
        content = '<p style="word-wrap: break-word; max-width: 595px; font-size: 14pt;">{}</p>'.format(
            content)
    filepath = f'{output_dir}/{filename}.pdf'
    # Create a PDF of the email and save it under the new directory
    HTML(string=content).write_pdf(filepath)
    return filepath


def process_directory_files(input_dir, output_dir):
    # Iterate over all subdirectories of input_dir
    for root, dirs, files in os.walk(input_dir):
        for dir in dirs:
            dir_path = os.path.join(root, dir)

            # Read the JSON file
            with open(os.path.join(dir_path, "info.json"), 'r') as f:
                info = json.load(f)

            sender = info['From']
            sender_dir, attachments_dir, text_dir, attachments_img_dir, text_img_dir = create_directories(
                output_dir, sender)

            # copy the json file to the sender directory
            shutil.copy(os.path.join(dir_path, "info.json"), sender_dir)

            attachment_filepaths = []

            # Iterate over all PDF files in the directory
            for file in os.listdir(dir_path):
                if file.endswith(".pdf"):
                    # Save the filepath
                    attachment_filepaths.append(os.path.join(dir_path, file))

            attachment_images = convert_pdfs_to_images(
                attachment_filepaths, attachments_img_dir)

            for image in attachment_images:

                image_name = get_image_name(image)
                sender_name = get_sender_name(sender)

                if not is_duplicate_image(image_name):
                    add_image_to_presentation(
                        image, image_name, sender_name)

            return


def process_eml_files(input_dir, output_dir):
    for filename in os.listdir(input_dir):
        if filename.endswith('.eml'):
            process_single_eml_file(filename, output_dir)
    return


def process_single_eml_file(filename, output_dir):
    with open(f'./eml/{filename}', 'rb') as f:
        msg = BytesParser(policy=policy.default).parse(f)

    sender = parseaddr(msg['From'])[1]
    sender_dir, attachments_dir, text_dir, attachments_img_dir, text_img_dir = create_directories(
        output_dir, sender)

    # copy the eml file to the sender directory
    shutil.copy(f'./eml/{filename}', sender_dir)

    message_content = get_html_content(msg)

    if message_content is None:
        raise Exception(
            f'No HTML/Text content found in {filename} for sender {sender}')

    filename = sanitize_filename(filename)

    text_filepath = create_pdf(message_content, text_dir, filename)
    text_images = convert_pdfs_to_images([text_filepath], text_img_dir)

    attachment_filepaths = save_attachments(msg, attachments_dir)

    attachment_images = convert_pdfs_to_images(
        attachment_filepaths, attachments_img_dir)

    all_images = text_images + attachment_images

    for image in all_images:

        image_name = get_image_name(image)
        sender_name = get_sender_name(sender)

        if not is_duplicate_image(image_name):
            add_image_to_presentation(
                image, image_name, sender_name)

    return


def is_duplicate_image(hash_image_name):
    return any(read_config_from_text_frame(slide, "image_name") == hash_image_name for slide in prs.slides)


def add_image_to_presentation(image, hashed_image_name, hashed_sender_name):
    # Check if the sender already exists in the dictionary
    if hashed_sender_name in slides_dict:
        # Add the image to the sender's list
        slides_dict[hashed_sender_name].append(image)
    else:
        # If the sender does not exist, create a new list for them
        slides_dict[hashed_sender_name] = [image]


def create_presentation_from_dict(prs):
    i = 0
    for sender, images in slides_dict.items():
        for image in images:
            i += 1
            image_name = get_image_name(image)
            # Add a new slide at the end
            # print(f'Adding slide nr. {i} {image} to slide with hashes {hashed_image_name}')
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            write_default_json_config_to_text_frame(slide)
            write_config_to_text_frame(
                slide, "image_name", image_name)
            write_config_to_text_frame(slide, "sender_name", sender)
            remove_title_placeholder(slide)
            add_header_left(slide)
            add_image(slide, image)
            text = f'Quelle: {sender} - {image_basename(image)}'
            if len(text) > 120:
                space_index = text.rfind(' ', 0, 120)
                if space_index != -1:
                    text = text[:space_index] + "\n" + text[space_index+1:]

            add_source_text(slide, text)
            add_text_box(slide)
            add_divider_line(slide, prs)


def add_headers_and_print_hashes(prs):
    page_count = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        # if right header exists, remove it
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if regex_right_header.match(text_frame.text):
                    slide.shapes._spTree.remove(shape._element)
        add_header_right(slide, i+1, page_count)
    return page_count

    # for i, slide in enumerate(prs.slides):
    #     # hashed_image_name = read_from_slide_note(slide,"hashed_image_name")
    #     # hashed_sender_name = read_from_slide_note(slide,"hashed_sender_name")
    #     hashed_image_name = read_config_from_text_frame(slide, "image_name")
    #     hashed_sender_name = read_config_from_text_frame(slide, "sender_name")
    #     print(f'{i+1}: {hashed_image_name} - {hashed_sender_name}')


def save_presentation(prs):
    prs.save(presentation_filename)


def group_consecutive_numbers(numbers):
    ranges = []
    for n in numbers:
        if not ranges or n > ranges[-1][-1] + 1:
            ranges.append([n])
        else:
            ranges[-1].append(n)
    return ['{}-{}'.format(r[0], r[-1]) if len(r) > 1 else str(r[0]) for r in ranges]


def get_all_senders(prs):
    all_senders = []
    for i, slide in enumerate(prs.slides):
        hashed_image_name = read_config_from_text_frame(slide, "image_name")
        hashed_sender_name = read_config_from_text_frame(slide, "sender_name")
        all_senders.append(hashed_sender_name)
    return all_senders


def get_sender_positions(all_senders):
    senders = [x for i, x in enumerate(
        all_senders) if i == 0 or x != all_senders[i-1]]
    duplicates = [item for item in senders if senders.count(item) > 1]
    senders = list(dict.fromkeys(duplicates))

    sender_str_list = []
    for sender in senders:
        positions = [i+1 for i, x in enumerate(all_senders) if x == sender]
        grouped_positions = group_consecutive_numbers(positions)
        sender_str_list.append(
            f'{sender} appears at pages {", ".join(grouped_positions)}')
    return sender_str_list


# Usage
if eml_input_dir:
    process_eml_files(eml_input_dir, output_dir)
if scanned_input_dir:
    process_directory_files(scanned_input_dir, output_dir)
create_presentation_from_dict(prs)
page_count = add_headers_and_print_hashes(prs)
save_presentation(prs)

# print that new pages start form page
if page_count > starting_page_count:
    print(f'New pages start from page: {page_count}')


all_senders = get_all_senders(prs)
sender_positions = get_sender_positions(all_senders)
if sender_positions:
    print("The following senders appear in more than one section in the presentation:")
    print("\n".join(sender_positions))
