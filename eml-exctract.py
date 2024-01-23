import os
import email
from email import policy
from email.parser import BytesParser
from weasyprint import HTML
from pdf2image import convert_from_path
from PIL import ImageDraw, ImageFont, Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.util import Mm, Cm
# from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import PP_PLACEHOLDER
import re
import hashlib
import json

text_types = ['text/plain', 'text/html']
input_dir = './eml'
output_dir = './output'
font_path = "/Library/Fonts/Arial.ttf"
page_count = 0

# Create output directory if it doesn't exist
if not os.path.exists(output_dir):
    os.makedirs(output_dir)


# Check if the PowerPoint file already exists
if os.path.exists('presentation.pptx'):
    prs = Presentation('presentation.pptx')
else:
    prs = Presentation()

# Set slide width and height to A4 portrait dimensions (210mm x 297mm)
prs.slide_width = Mm(297)
prs.slide_height = Mm(210)

def image_basename(image):
    return os.path.splitext(os.path.basename(image))[0]

def hash_image_name(image_name):
    name = image_basename(image_name)
    hasher = hashlib.md5()
    hasher.update(name.encode('utf-8'))
    unique_representation = hasher.hexdigest()
    return "hash_image_" + unique_representation



def hash_sender_name(sender):  
    hasher = hashlib.md5()
    hasher.update(sender.encode('utf-8'))
    unique_representation = hasher.hexdigest()
    return "hash_sender_" + unique_representation

def write_default_json_to_notes(slide):
    default_for_json = {"hashed_image_name": "default", "hashed_sender_name": "default"}
    json_hashed_image_name = json.dumps(default_for_json)
    slide.notes_slide.notes_text_frame.text = json_hashed_image_name

def write_to_slide_note(slide, key, value):
    # Read the existing content
    existing_content = slide.notes_slide.notes_text_frame.text
    existing_json = json.loads(existing_content) if existing_content else {}

    # Update the JSON
    existing_json[key] = value

    # Write it back
    slide.notes_slide.notes_text_frame.text = json.dumps(existing_json)

def read_from_slide_note(slide, key):
    # Read the existing content
    existing_content = slide.notes_slide.notes_text_frame.text
    existing_json = json.loads(existing_content) if existing_content else {}

    # Update the JSON
    return existing_json[key]

def sanitize_filename(filename):
    return re.sub(r'[\\/*?:"<>|]', '_', filename)

# conda install conda-forge::weasyprint 

def get_html_content(msg):
    if msg.is_multipart():
        for part in msg.iter_parts():
            # Print content_type
            if part.get_content_type() in text_types:
                return part.get_content()
            elif part.is_multipart():
                return get_html_content(part)
    else:
        if msg.get_content_type() == 'text/html':
            return msg.get_content()

    return None

from PIL import Image
from PIL import ImageChops

def add_header_left(slide):
    bold = False
    size = 11
    # Add a header to the slide
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = "Abwägungsvorschlag Träger öffentlicher Belange, B-Plan XX “XXXX“"
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_header_right(slide, page_number, total_pages=page_count):
    bold = False
    size = 10
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = f'Seite {page_number}/{total_pages}'
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

def crop_whitespace(image):
    # Crop the image to remove whitespace
    bg = Image.new(image.mode, image.size, image.getpixel((0,0)))
    diff = ImageChops.difference(image, bg)
    diff = ImageChops.add(diff, diff, 2.0, -100)
    bbox = diff.getbbox()
    if bbox:
        return image.crop(bbox)

def add_image(slide, image_filepath):
    # Open the image and get its size
    image = Image.open(image_filepath)
    image_width, image_height = image.size
    aspect_ratio = image_width / image_height

    # Define the maximum height and width
    max_height = Cm(18)
    max_width = Cm(13)

    # Calculate the height and width, maintaining the aspect ratio
    if aspect_ratio > 1:
        # Image is wider than it is tall
        width = max_width
        height = width / aspect_ratio
    else:
        # Image is taller than it is wide
        height = max_height
        width = height * aspect_ratio

    # Add the image to the slide
    top = Cm(1.6)
    left = Cm(1)
    pic = slide.shapes.add_picture(image_filepath, left, top, height=height, width=width)

    # Crop the image
    pic.crop_left = 0
    pic.crop_top = 0
    pic.crop_right = 0
    pic.crop_bottom = 0

# Function that adds a string to the top left of the slide
def add_text(slide, text):
    left = top = width = height = Cm(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(6)

def add_divider_line(slide, prs):
    left = prs.slide_width // 2
    top = Cm(0.7)
    end_top = prs.slide_height
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left, end_top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False

    # Add horizontal divider line 1 cm from the top
    left = Cm(0.5)
    top = Cm(0.7)
    end_left = prs.slide_width - Cm(0.5)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, end_left, top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False


def add_text_box(slide):
    left = Cm(16)
    top = Cm(1)
    width = Cm(11)
    height = Cm(20)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Zur Kenntnis genommen."
    p.font.bold = False
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.LEFT

def remove_title_placeholder(slide):
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == 0 and phf.type == PP_PLACEHOLDER.TITLE:
                sp = shape._element
                sp.getparent().remove(sp)

def create_directories(output_dir, sender):
    # Create a new directory with the same name as the sender's email address
    sender_dir = f'{output_dir}/{sender}'
    os.makedirs(sender_dir, exist_ok=True)
    
    # Create a subdirectory called "attachments" to store attachments
    attachments_dir = f'{sender_dir}/attachments'
    os.makedirs(attachments_dir, exist_ok=True)
    
    # Create a subdirectory called "text" to store HTML files
    text_dir = f'{sender_dir}/text'
    os.makedirs(text_dir, exist_ok=True)
    
    # Create a subdirectory called "img" in subdirectory "attachments" to store images
    attachments_img_dir = f'{attachments_dir}/img'
    os.makedirs(attachments_img_dir, exist_ok=True)
    
    # Create a subdirectory called "img" in subdirectory "text" to store images
    text_img_dir = f'{text_dir}/img'
    os.makedirs(text_img_dir, exist_ok=True)

    return sender_dir, attachments_dir, text_dir, attachments_img_dir, text_img_dir
    

def save_attachments(msg, output_dir):
    filepaths = []
    if msg.is_multipart():
        for part in msg.iter_parts():
            if part.get_content_type() not in text_types:
                payload = part.get_payload(decode=True)
                if payload is not None:
                    attachment_filename = part.get_filename()
                    filepath = f'{output_dir}/{attachment_filename}'
                    with open(filepath, 'wb') as f:
                        f.write(payload)
                    filepaths.append(filepath)
    return filepaths

def convert_pdfs_to_images(pdf_filepaths, output_dir):
    images = []
    for pdf_filepath in pdf_filepaths:
        # Convert the PDF to images
        pdf_images = convert_from_path(pdf_filepath)
        # Save each image under the output directory
        for i, pdf_image in enumerate(pdf_images):
            pdf_image = crop_whitespace(pdf_image)
            image_filepath = f'{output_dir}/{os.path.basename(pdf_filepath)}_{i}.png'
            pdf_image.save(image_filepath, 'PNG')
            images.append(image_filepath)
    return images

def create_pdf(html_content, output_dir, filename):
    filepath = f'{output_dir}/{filename}.pdf'
    # Create a PDF of the email and save it under the new directory
    HTML(string=html_content).write_pdf(filepath)
    return filepath

def process_eml_files(input_dir, output_dir):
    page_count = 0
    for filename in os.listdir(input_dir):
        if filename.endswith('.eml'):
            page_count = process_single_eml_file(filename, output_dir, page_count)
    return page_count

def process_single_eml_file(filename, output_dir, page_count):
    with open(f'./eml/{filename}', 'rb') as f:
        msg = BytesParser(policy=policy.default).parse(f)
    
    sender = email.utils.parseaddr(msg['From'])[1]
    sender_dir, attachments_dir, text_dir, attachments_img_dir, text_img_dir = create_directories(output_dir, sender)

    html_content = get_html_content(msg)
    if html_content is None:
        print(f'No HTML content found in {filename}')
    filename = sanitize_filename(filename)

    text_filepath = create_pdf(html_content, text_dir, filename)
    text_images = convert_pdfs_to_images([text_filepath], text_img_dir)
    
    attachment_filepaths = save_attachments(msg, attachments_dir)
    print(attachment_filepaths)
                    
    attachment_images = convert_pdfs_to_images(attachment_filepaths, attachments_img_dir)

    all_images = text_images + attachment_images

    for image in all_images:
        page_count += 1
        print(f'Page count {page_count}')
        
        hashed_image_name = hash_image_name(image)
        hashed_sender_name = hash_sender_name(sender)

        if not is_duplicate_image(hashed_image_name):
            add_image_to_presentation(image, hashed_image_name, hashed_sender_name, page_count)

    print("------------------------------------------")
    return page_count

def is_duplicate_image(hashed_image_name):
    return any(read_from_slide_note(slide, "hashed_image_name") == hashed_image_name for slide in prs.slides)

def add_image_to_presentation(image, hashed_image_name, hashed_sender_name, page_count):
    # Check if the sender already exists in the presentation
    sender_exists = any(read_from_slide_note(slide, "hashed_sender_name") == hashed_sender_name for slide in prs.slides)

    if sender_exists:
        # Find the index of the last slide of this sender
        last_slide_index = max(i for i, slide in enumerate(prs.slides) if read_from_slide_note(slide, "hashed_sender_name") == hashed_sender_name)

        # Insert a new slide after the last slide of this sender
        slide = prs.slides.add_slide(prs.slide_layouts[5], last_slide_index + 1)
    else:
        # If the sender does not exist, add a new slide at the end
        slide = prs.slides.add_slide(prs.slide_layouts[5])

    remove_title_placeholder(slide)
    print(f'Adding {image} to slide {page_count} with hashes {hashed_image_name}')
    write_to_slide_note(slide, "hashed_image_name", hashed_image_name)
    write_to_slide_note(slide, "hashed_sender_name", hashed_sender_name)
    add_header_left(slide)
    add_image(slide, image)
    add_text(slide, f'{image_basename(image)}')
    add_text_box(slide)
    add_divider_line(slide, prs)

def add_headers_and_print_hashes(prs, page_count):
    for i, slide in enumerate(prs.slides):
        add_header_right(slide, i+1, page_count)

    for i, slide in enumerate(prs.slides):
        hashed_image_name = read_from_slide_note(slide,"hashed_image_name")
        hashed_sender_name = read_from_slide_note(slide,"hashed_sender_name")
        print(f'{i+1}: {hashed_image_name} - {hashed_sender_name}')
        # print(f'{i+1}: {hashed_image_name}')

def save_presentation(prs):
    prs.save('presentation.pptx')

# Usage
page_count = process_eml_files(input_dir, output_dir)
add_headers_and_print_hashes(prs, page_count)
save_presentation(prs)