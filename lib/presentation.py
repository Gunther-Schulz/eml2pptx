import os
import re
from pptx import Presentation
from pptx import Presentation
from pptx.util import Pt, Mm, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, PP_PLACEHOLDER
from PIL import Image

from lib.config_manager import get_processed_slides_from_json_file, get_slide_id, image_basename, load_config, read_config_from_text_box, write_config_to_text_box, write_default_json_config_to_text_frame, write_processed_slide_to_json_file


config = load_config()
presentation_filename = config['presentation_filename']
header_title = config['header_title']
default_comment = config['default_comment']

slides_dict = {}
regex_right_header = re.compile(r"Seite \d+/\d+")


# Check if the PowerPoint file already exists
if os.path.exists(presentation_filename):
    # If it does, open it
    prs = Presentation(presentation_filename)
    starting_page_count = len(prs.slides)

else:
    prs = Presentation()

prs.slide_width = Mm(297)
prs.slide_height = Mm(210)


def is_duplicate(image, sender):
    return any((read_config_from_text_box(slide, "image") == image and read_config_from_text_box(slide, "sender") == sender) for slide in prs.slides)


def add_header_left(slide):
    bold = False
    size = 11
    # Add a header to the slide
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = header_title
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT


def add_header_right(slide, page_number, total_pages):
    bold = False
    size = 10
    header = slide.shapes.add_textbox(Cm(1), 0, prs.slide_width - Cm(2), Cm(1))
    tf = header.text_frame
    tf.text = f'Seite {page_number}/{total_pages}'
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_image(slide, image_filepath):
    # Define the maximum height and width
    max_height = Cm(18)
    max_width = Cm(14)

    # Open the image and get its size
    image = Image.open(image_filepath)
    image_width, image_height = image.size
    aspect_ratio = image_width / image_height

    # Calculate the height and width, maintaining the aspect ratio
    width_by_height = max_height * aspect_ratio
    height_by_width = max_width / aspect_ratio

    if width_by_height <= max_width:
        # Image fits within max dimensions by adjusting height
        width = width_by_height
        height = max_height
    else:
        # Image fits within max dimensions by adjusting width
        width = max_width
        height = height_by_width

    # Add the image to the slide
    top = Cm(1.6)
    left = Cm(0.5)
    pic = slide.shapes.add_picture(
        image_filepath, left, top, height=height, width=width)


def add_source_text(slide, text):
    # Function that adds a string to the top left of the slide
    left = Cm(1.6)
    top = Cm(20)
    width = Cm(13.5)
    height = Cm(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(6)


def add_divider_line(slide, prs):
    left = prs.slide_width // 2
    top = Cm(0.7)
    end_top = prs.slide_height
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, left, top, left, end_top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False

    # Add horizontal divider line
    left = Cm(0.5)
    top = Cm(0.7)
    end_left = prs.slide_width - Cm(0.5)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, left, top, end_left, top)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(1)
    line.shadow.inherit = False
    line.shadow.visible = False


def add_text_box(slide):
    left = Cm(16)
    top = Cm(1)
    width = Cm(11)
    height = Cm(20)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = default_comment
    p.font.bold = False
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.LEFT


def remove_title_placeholder(slide):
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            if phf.idx == 0 and phf.type == PP_PLACEHOLDER.TITLE:
                sp = shape._element
                sp.getparent().remove(sp)


def add_image_to_presentation(image, sender):
    slide_id = get_slide_id(image, sender)
    if slide_id not in get_processed_slides_from_json_file():
        if sender in slides_dict:
            # Add the image to the sender's list
            slides_dict[sender].append(image)
        else:
            # If the sender does not exist, create a new list for them
            slides_dict[sender] = [image]
    else:
        print(f'Slide was removed manually. Not re-adding. {slide_id}')
    write_processed_slide_to_json_file(slide_id)


def create_presentation_from_dict():
    i = 0
    for sender, images in slides_dict.items():
        for image in images:
            i += 1
            # Add a new slide at the end
            # print(f'Adding slide nr. {i} {image} to slide with hashes {hashed_image_name}')
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            write_default_json_config_to_text_frame(slide)
            write_config_to_text_box(
                slide, "image", image)
            write_config_to_text_box(slide, "sender", sender)
            remove_title_placeholder(slide)
            add_header_left(slide)
            add_image(slide, image)
            text = f'Quelle: {sender} - {image_basename(image)}'
            if len(text) > 120:
                space_index = text.rfind(' ', 0, 120)
                if space_index != -1:
                    text = text[:space_index] + "\n" + text[space_index+1:]

            add_source_text(slide, text)
            add_text_box(slide)
            add_divider_line(slide, prs)
    return prs


def add_headers(prs):
    page_count = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        # if right header exists, remove it
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if regex_right_header.match(text_frame.text):
                    slide.shapes._spTree.remove(shape._element)
        add_header_right(slide, i+1, page_count)
    return page_count


def save_presentation(prs):
    prs.save(presentation_filename)
