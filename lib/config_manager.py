import json
import os
from email.utils import parsedate_to_datetime
import re
import yaml
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor


def load_config():
    if not os.path.exists('config.yaml'):
        print("config.yaml does not exist. Creating it with default values.")
        with open('config.yaml', 'w') as f:
            f.write('presentation_filename: presentation.pptx\n')
            f.write('header_title: "Eingegangene Emails"\n')
            f.write('default_comment: "Wird zur Kenntnis genommen."\n')
            f.write('pdf_blacklist:\n')
            f.write('  - ".*DUMMY.*"\n')
            f.write('eml_input_dir: eml\n')
            f.write('scanned_input_dir: scanned\n')
            f.write('output_dir: extracted\n')
            # f.write('font_path: "/Library/Fonts/Arial.ttf"\n')
        print("Please configure config.yaml and run the script again.")
        exit()

    with open('config.yaml', 'r') as f:
        config = yaml.safe_load(f)

    config['pdf_blacklist'] = [re.compile(
        pattern) for pattern in config['pdf_blacklist']]
    return config


config = load_config()
pdf_blacklist = config['pdf_blacklist']
output_dir = config['output_dir']

processed_slides_file = 'processed_slides.json'

# Create output directory if it doesn't exist
if not os.path.exists(output_dir):
    os.makedirs(output_dir)


def image_basename(image):
    return os.path.splitext(os.path.basename(image))[0]


if not os.path.exists(processed_slides_file):
    with open(processed_slides_file, 'w') as f:
        json.dump({'processed_slides': []}, f)


def create_directories(output_dir, sender):
    # Create a new directory with the same name as the sender's email address
    sender_dir = f'{output_dir}/{sender}'
    os.makedirs(sender_dir, exist_ok=True)

    # Create a subdirectory called "attachments" to store attachments
    attachments_dir = f'{sender_dir}/attachments'
    os.makedirs(attachments_dir, exist_ok=True)

    # Create a subdirectory called "text" to store HTML files
    text_dir = f'{sender_dir}/text'
    os.makedirs(text_dir, exist_ok=True)

    # Create a subdirectory called "img" in subdirectory "attachments" to store images
    attachments_img_dir = f'{attachments_dir}/img'
    os.makedirs(attachments_img_dir, exist_ok=True)

    # Create a subdirectory called "img" in subdirectory "text" to store images
    text_img_dir = f'{text_dir}/img'
    os.makedirs(text_img_dir, exist_ok=True)

    return sender_dir, attachments_dir, text_dir, attachments_img_dir, text_img_dir


def get_email_date(msg):
    email_date = parsedate_to_datetime(msg['Date'])
    # Format the date and time as YYYY-MM-DD_HH_MM_SS
    formatted_date_time = email_date.strftime('%Y-%m-%d_%H_%M_%S')
    return formatted_date_time


def sanitize_filename(filename):
    return re.sub(r'[\\/*?:"<>|\[\]]', '_', filename)


def is_in_blacklist(filename):
    for regex in pdf_blacklist:
        if regex.match(filename):
            print(f'Found blacklisted file: {filename}')
            return True
    return False


def get_processed_slides_from_json_file():
    with open(processed_slides_file) as f:
        data = json.load(f)
        return data['processed_slides']


def write_processed_slide_to_json_file(slide_id):
    with open(processed_slides_file, 'r+') as f:
        data = json.load(f)
        data['processed_slides'].append(slide_id)
        f.seek(0)
        json.dump(data, f)


def get_slide_id(image, sender):
    # image_name = get_image_name(image)
    # sender_name = get_sender_name(sender)
    # slide_id = hashlib.md5(
    #     (image_name + sender_name).encode('utf-8')).hexdigest()
    slide_id = image + "_&_" + sender
    return slide_id


def write_default_json_config_to_text_frame(slide):
    # create a new text box, not a note
    text_frame = slide.shapes.add_textbox(Cm(0.1), Cm(0), Cm(0.1), Cm(
        0.1)).text_frame  # Parameters are left, top, width, height
    # write default json to it
    default_for_json = {"image": "default",
                        "sender": "default"}
    json_image = json.dumps(default_for_json)
    text_frame.text = json_image
    text_frame.paragraphs[0].font.size = Pt(1)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def get_json_text_box(slide):
    # search all text boxes and find the first one that contains json
    regex_json = re.compile(r"^{.*}$")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if regex_json.match(text_frame.text):
                return text_frame
    return None


def write_config_to_text_box(slide, key, value):
    json_text_box = get_json_text_box(slide)
    if json_text_box is None:
        return
    # get the json from the text box and update it
    existing_json = json.loads(json_text_box.text)
    existing_json[key] = value
    # write the updated json back to the text box
    json_text_box.text = json.dumps(existing_json)
    json_text_box.paragraphs[0].font.size = Pt(1)
    json_text_box.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def read_config_from_text_box(slide, key):
    json_text_box = get_json_text_box(slide)
    if json_text_box is None:
        return None
    # get the json from the text box
    existing_json = json.loads(json_text_box.text)
    return existing_json.get(key, None)
